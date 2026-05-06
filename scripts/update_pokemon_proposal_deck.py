from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "pokemon_ecs163_proposal_fixed_layout.pptx"
OUTPUT = ROOT / "pokemon_ecs163_proposal_final_with_assets.pptx"
ASSETS = ROOT / "data" / "presentation_assets"

IMAGE_BOX = (7525512, 1764792, 3739896, 3008376)
TEXT_COLOR = RGBColor(255, 247, 235)
MUTED_COLOR = RGBColor(184, 192, 200)


def clone_slide(prs, source_slide):
    blank = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank)
    for shape in source_slide.shapes:
        new_slide.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")
    return new_slide


def move_slide(prs, old_index, new_index):
    slide_id = prs.slides._sldIdLst[old_index]
    prs.slides._sldIdLst.remove(slide_id)
    prs.slides._sldIdLst.insert(new_index, slide_id)


def delete_shape(shape):
    shape.element.getparent().remove(shape.element)


def set_text(shape, text, size=None, bold=None, color=None):
    shape.text = text
    for paragraph in shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            if size:
                run.font.size = Pt(size)
            if bold is not None:
                run.font.bold = bold
            if color:
                run.font.color.rgb = color


def text_shapes(slide):
    return [shape for shape in slide.shapes if hasattr(shape, "text")]


def find_shape(slide, startswith=None, contains=None):
    for shape in text_shapes(slide):
        txt = shape.text.strip()
        if startswith and txt.startswith(startswith):
            return shape
        if contains and contains in txt:
            return shape
    return None


def clear_insert_placeholders(slide):
    for shape in list(text_shapes(slide)):
        if shape.text.strip().startswith("Insert"):
            delete_shape(shape)


def add_image(slide, image_name):
    x, y, w, h = IMAGE_BOX
    path = ASSETS / image_name
    slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def update_footer(slide, number):
    for shape in text_shapes(slide):
        if "ECS 163 Proposal" in shape.text:
            set_text(shape, f"ECS 163 Proposal  •  {number:02d}", size=9, color=MUTED_COLOR)


def set_bullets(slide, bullets):
    body = None
    for shape in text_shapes(slide):
        if shape.left < 2_000_000 and shape.top > 1_300_000 and shape.height > 2_000_000:
            body = shape
            break
    if body:
        set_text(body, "\n".join(bullets), size=18)

    # The source deck has a decorative card plus an older body text box.
    # Keep the card text and remove the duplicate so old copy cannot show through.
    for shape in list(text_shapes(slide)):
        if shape.left == 914400 and shape.top == 1783080 and shape.width == 5715000:
            delete_shape(shape)


def set_title_subtitle(slide, title, subtitle):
    title_shape = text_shapes(slide)[9]
    subtitle_shape = text_shapes(slide)[10]
    set_text(title_shape, title, size=31, bold=True, color=TEXT_COLOR)
    set_text(subtitle_shape, subtitle, size=15, color=MUTED_COLOR)


def set_caption(slide, caption):
    for shape in text_shapes(slide):
        txt = shape.text.strip()
        if txt in {"Dataset + findings", "Interesting insights", "Visualization idea", "Story plan"}:
            set_text(shape, caption, size=10, bold=True)
            return


def update_content_slide(slide, title, subtitle, bullets, image_name, caption=None):
    set_title_subtitle(slide, title, subtitle)
    set_bullets(slide, bullets)
    clear_insert_placeholders(slide)
    add_image(slide, image_name)
    if caption:
        set_caption(slide, caption)


def update_script_slide(slide, title, body):
    set_title_subtitle(slide, title, "Speaker notes — not for main presentation")
    for shape in text_shapes(slide):
        if shape.left > 1_000_000 and shape.top > 1_600_000 and shape.width > 8_000_000:
            set_text(shape, body, size=21)
            break


def main():
    prs = Presentation(SOURCE)

    # Add one missing content slide after insights, then one matching script slide after script 3.
    clone_slide(prs, prs.slides[3])
    move_slide(prs, len(prs.slides) - 1, 3)
    clone_slide(prs, prs.slides[10])
    move_slide(prs, len(prs.slides) - 1, 11)

    content = [
        (
            "What Really Makes a Pokémon Strong?",
            "Visualizing competitive success beyond raw stats",
            [
                "Common assumption: higher base stats mean stronger Pokémon",
                "Our claim: competitive success also depends on synergy",
                "Narrative visualization reveals the hidden team ecosystem",
            ],
            "network_mockup.png",
            None,
        ),
        (
            "Dataset and Research Focus",
            "Rubric: dataset description + interesting insights",
            [
                "Competitive data includes stats, usage, teammates, moves, items, abilities",
                "Cleaned outputs: Pokémon nodes, team edges, build usage rows, image lookup",
                "1,098 Pokémon/forms, 2,606 team edges, 2,967 build rows",
                "Focus: VGC team ecosystems, not a generic Pokédex browser",
            ],
            "dataset_pipeline.png",
            "Dataset pipeline",
        ),
        (
            "Interesting Insights",
            "The data challenges the “stats = strength” assumption",
            [
                "Total stats vs usage correlation: 0.194",
                "Team connectivity vs usage correlation: 0.903",
                "Incineroar: moderate stats, strongest connectivity signal",
                "High stats can still have low competitive usage",
            ],
            "correlation_insight.png",
            "Interesting insights",
        ),
        (
            "Case Studies: Strength Is Contextual",
            "Recognizable examples make the contradiction easy to explain",
            [
                "Zacian: high stats and high usage",
                "Zamazenta/Mewtwo: high stats but low usage",
                "Incineroar: moderate stats, very high team centrality",
                "Amoonguss/Grimmsnarl: support value comes through partners",
            ],
            "case_study_cards.png",
            "Case studies",
        ),
        (
            "Main Visualization Idea",
            "Rubric: visualization design and interaction",
            [
                "Advanced D3 force-directed synergy network",
                "Nodes = Pokémon; edges = common teammates",
                "Node size = usage percentage; edge thickness = co-usage strength",
                "Click a Pokémon to reveal teammates, builds, stats, and role",
            ],
            "network_mockup.png",
            "Visualization idea",
        ),
        (
            "Storytelling Structure: Martini Glass",
            "Rubric: how we communicate the story effectively",
            [
                "Start with a guided question: are raw stats enough?",
                "Reveal contradictions through annotated examples",
                "Transition into the team synergy network",
                "Open the visualization for user-driven exploration",
            ],
            "martini_glass_flow.png",
            "Story plan",
        ),
        (
            "Prototype and Next Steps",
            "Current status and final implementation plan",
            [
                "Working React + D3 prototype: hover, click, highlight, reset, detail panel",
                "Current demo focuses on a readable top competitive subset",
                "Next: guided transitions, annotations, clustering, and polish",
                "Final goal: a polished interactive story, not a dashboard",
            ],
            "network_mockup.png",
            None,
        ),
        (
            "Thank You",
            "Questions?",
            [
                "Project: What Really Makes a Pokémon Strong?",
                "Core thesis: competitive success emerges from synergy",
                "Final story will reveal why strong teams beat isolated stats",
            ],
            "network_blur.png",
            None,
        ),
    ]

    for idx, args in enumerate(content):
        update_content_slide(prs.slides[idx], *args)

    scripts = [
        (
            "Slide 1 Script — Opening Hook",
            "We start with a question that most Pokémon players understand: what actually makes a Pokémon strong? A beginner might assume the answer is simple: higher base stats. Our project challenges that assumption by showing that competitive strength also depends on team context.",
        ),
        (
            "Slide 2 Script — Dataset",
            "Our dataset comes from a competitive Pokémon database. We cleaned it into Pokémon nodes, teammate edges, build usage rows, and image references. That structure lets us study stats, usage, teammates, moves, items, and abilities without turning the project into a Pokédex browser.",
        ),
        (
            "Slide 3 Script — Insights",
            "The first insight is that raw stats only weakly explain competitive usage. Total stats and usage have a correlation of 0.194, while team connectivity and usage are much more closely related at 0.903. That supports our claim that competitive success is often relational.",
        ),
        (
            "Slide 4 Script — Case Studies",
            "The case studies make the pattern easier to understand. Zacian is strong by both stats and usage, but Zamazenta or Mewtwo show that high stats alone are not enough. Incineroar and support Pokémon like Amoonguss show how team roles can create competitive value.",
        ),
        (
            "Slide 5 Script — Visualization Idea",
            "Our main visualization is a force-directed team synergy network. Each node is a Pokémon, and each edge means two Pokémon commonly appear on teams together. Larger nodes show higher usage, thicker edges show stronger teammate relationships, and clicking a node reveals builds and partners.",
        ),
        (
            "Slide 6 Script — Storytelling Structure",
            "We chose a Martini Glass structure. The opening is author-driven: we introduce the raw-stat assumption, then reveal examples that contradict it. After that, the visualization opens into exploration, where users can inspect the team ecosystem themselves.",
        ),
        (
            "Slide 7 Script — Prototype and Future Work",
            "We already have a working React and D3 prototype with hover tooltips, click selection, teammate highlighting, reset behavior, and a detail panel. For the final project, we will add guided annotations, transitions, clustering, and polish so the experience feels like an interactive story rather than a dashboard.",
        ),
        (
            "Slide 8 Script — Closing",
            "To summarize, our project asks what really makes a Pokémon strong. Our answer is that strength is not only visible in base stats; it also emerges from hidden relationships between teammates, builds, and strategy. Thank you, and we are happy to take questions.",
        ),
    ]

    for offset, args in enumerate(scripts, start=8):
        update_script_slide(prs.slides[offset], *args)

    for i, slide in enumerate(prs.slides, start=1):
        update_footer(slide, i)

    prs.save(OUTPUT)
    print(f"Saved {OUTPUT}")


if __name__ == "__main__":
    main()
